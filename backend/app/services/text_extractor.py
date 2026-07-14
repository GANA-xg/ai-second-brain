"""
Text extraction from supported document formats.

Extracts structured text from PDF, DOCX, PPTX, TXT files.
Image files are detected and reported as unsupported (OCR not yet implemented).

Each extractor returns a list of ExtractionResult named tuples with:
  - text: extracted text content
  - page_number: page number (PDF only, 1-indexed)
  - slide_number: slide number (PPTX only, 1-indexed)
  - section: section heading (DOCX only, when available)
"""

from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional


@dataclass
class ExtractionResult:
    text: str
    page_number: Optional[int] = None
    slide_number: Optional[int] = None
    section: Optional[str] = None


# --------------------------------------------------------------------------- #
# PDF
# --------------------------------------------------------------------------- #


def extract_pdf(filepath: Path, max_pages: int = 100) -> List[ExtractionResult]:
    """Extract text from a PDF file, preserving page boundaries and reading order.

    Args:
        filepath: Path to the PDF file.
        max_pages: Maximum number of pages to extract (default 100).

    Returns:
        A list of ExtractionResult, one per page.

    Raises:
        ValueError: If the PDF cannot be read or contains no pages.
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        raise ImportError("pypdf is required for PDF extraction: pip install pypdf")

    try:
        reader = PdfReader(str(filepath))
    except Exception as exc:
        raise ValueError(f"Failed to read PDF: {exc}") from exc

    pages = reader.pages
    if not pages:
        raise ValueError("PDF file contains no pages")

    results: List[ExtractionResult] = []
    for i, page in enumerate(pages, start=1):
        if i > max_pages:
            break
        try:
            text = page.extract_text() or ""
        except Exception as exc:
            raise ValueError(f"Failed to extract text from page {i}: {exc}") from exc

        results.append(ExtractionResult(text=text, page_number=i))

    if not results:
        raise ValueError("No text could be extracted from PDF")

    return results


# --------------------------------------------------------------------------- #
# DOCX
# --------------------------------------------------------------------------- #


def extract_docx(filepath: Path, max_pages: int = 100, **kwargs: object) -> List[ExtractionResult]:
    """Extract text from a DOCX file, preserving paragraph boundaries and headings.

    Args:
        filepath: Path to the DOCX file.

    Returns:
        A list of ExtractionResult, one per paragraph.

    Raises:
        ValueError: If the DOCX cannot be read.
    """
    try:
        from docx import Document as DocxDocument
    except ImportError:
        raise ImportError(
            "python-docx is required for DOCX extraction: pip install python-docx"
        )

    try:
        doc = DocxDocument(str(filepath))
    except Exception as exc:
        raise ValueError(f"Failed to read DOCX: {exc}") from exc

    results: List[ExtractionResult] = []
    current_section: Optional[str] = None

    for para in doc.paragraphs:
        text = para.text
        if not text.strip():
            # Preserve paragraph structure with an empty line marker
            results.append(ExtractionResult(text="", section=current_section))
            continue

        # Detect section headings
        style_name = para.style.name if para.style else ""
        if style_name.startswith("Heading"):
            current_section = text

        results.append(ExtractionResult(text=text, section=current_section))

    if not results:
        # Truly empty document
        return []

    return results


# --------------------------------------------------------------------------- #
# PPTX
# --------------------------------------------------------------------------- #


def extract_pptx(filepath: Path, **kwargs: object) -> List[ExtractionResult]:
    """Extract text from a PPTX file, slide by slide.

    Args:
        filepath: Path to the PPTX file.

    Returns:
        A list of ExtractionResult, one per slide.

    Raises:
        ValueError: If the PPTX cannot be read or contains no slides.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx is required for PPTX extraction: pip install python-pptx"
        )

    try:
        prs = Presentation(str(filepath))
    except Exception as exc:
        raise ValueError(f"Failed to read PPTX: {exc}") from exc

    slides = prs.slides
    if not slides:
        raise ValueError("PPTX file contains no slides")

    results: List[ExtractionResult] = []
    for i, slide in enumerate(slides, start=1):
        slide_text: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)

        combined = "\n".join(slide_text)
        results.append(ExtractionResult(text=combined, slide_number=i))

    return results


# --------------------------------------------------------------------------- #
# TXT
# --------------------------------------------------------------------------- #


def extract_txt(filepath: Path, **kwargs: object) -> List[ExtractionResult]:
    """Extract text from a plain UTF-8 text file.

    Args:
        filepath: Path to the TXT file.

    Returns:
        A single-element list of ExtractionResult containing the full text.

    Raises:
        ValueError: If the file cannot be read.
    """
    try:
        raw = filepath.read_bytes()
        text = raw.decode("utf-8", errors="replace")
    except Exception as exc:
        raise ValueError(f"Failed to read TXT file: {exc}") from exc

    # Normalize line endings at extraction time
    text = text.replace("\r\n", "\n").replace("\r", "\n")

    return [ExtractionResult(text=text)]


# --------------------------------------------------------------------------- #
# Image detection
# --------------------------------------------------------------------------- #


def extract_image(filepath: Path, **kwargs: object) -> List[ExtractionResult]:
    """Handle image uploads — OCR is not yet implemented.

    Args:
        filepath: Path to the image file (PNG, JPG, JPEG).

    Returns:
        Never returns — always raises ValueError.

    Raises:
        ValueError: With a clear message that OCR is not yet implemented.
    """
    raise ValueError(
        "OCR text extraction is not yet implemented for image files. "
        "Image uploads are accepted but cannot be processed for text content "
        "in this phase."
    )


# --------------------------------------------------------------------------- #
# Dispatch
# --------------------------------------------------------------------------- #

# Map extension (without dot, lowercase) to extractor function
EXTRACTORS = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "pptx": extract_pptx,
    "txt": extract_txt,
    "png": extract_image,
    "jpg": extract_image,
    "jpeg": extract_image,
}


def get_extractor(extension: str):
    """Get the extractor function for a file extension.

    Args:
        extension: File extension (without dot, lowercase).

    Returns:
        The extractor function for this extension.

    Raises:
        ValueError: If the extension is not supported.
    """
    ext = extension.lower()
    if ext not in EXTRACTORS:
        raise ValueError(f"Unsupported file extension for extraction: '{ext}'")
    return EXTRACTORS[ext]
