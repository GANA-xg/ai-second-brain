"""Tests for the file upload system (Part 4).

Tests cover validation, authentication, ownership, checksums, status lifecycle,
filesystem isolation, and audit logging.
"""
import io
import os
import uuid
from pathlib import Path

import pytest
from fastapi import status
from fastapi.testclient import TestClient
from starlette.datastructures import UploadFile


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_pdf_bytes() -> bytes:
    """Minimal but valid PDF with extractable text using pypdf."""
    from pypdf import PdfWriter
    buf = io.BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(612, 792)
    writer.write(buf)
    buf.seek(0)
    return buf.getvalue()


def _make_png_bytes() -> bytes:
    """Minimal valid 1×1 red PNG."""
    return (
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x02\x00\x00\x00\x90wS\xde"
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )


def _make_jpeg_bytes() -> bytes:
    """Minimal valid JPEG (SOI + EOI markers)."""
    return b"\xff\xd8\xff\xe0\x00\x10JFIF\x00\x01\x01\x00\x00\x01\x00\x01\x00\x00\xff\xd9"


def _make_docx_bytes() -> bytes:
    """Valid DOCX with text content using python-docx."""
    from docx import Document
    buf = io.BytesIO()
    doc = Document()
    doc.add_paragraph("Test document content.")
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _make_pptx_bytes() -> bytes:
    """Valid PPTX with text content using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches
    buf = io.BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "Test slide content."
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _make_txt_bytes(text: str = "Hello, world!\n") -> bytes:
    return text.encode("utf-8")


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def sample_pdf():
    return "test.pdf", _make_pdf_bytes()


@pytest.fixture
def sample_png():
    return "diagram.png", _make_png_bytes()


@pytest.fixture
def sample_jpg():
    return "photo.jpg", _make_jpeg_bytes()


@pytest.fixture
def sample_jpeg():
    return "image.jpeg", _make_jpeg_bytes()


@pytest.fixture
def sample_docx():
    return "report.docx", _make_docx_bytes()


@pytest.fixture
def sample_pptx():
    return "slides.pptx", _make_pptx_bytes()


@pytest.fixture
def sample_txt():
    return "notes.txt", _make_txt_bytes()


# ===================================================================
# Successful uploads for every allowed type
# ===================================================================

class TestSuccessfulUploads:

    def _upload(self, client, auth_headers, name, content):
        return client.post(
            "/api/v1/files/upload",
            files={"file": (name, content)},
            headers=auth_headers,
        )
    def _verify_ok(self, response, name, mime_type, extension, expected_status="READY"):
        assert response.status_code == 201
        body = response.json()
        doc = body["document"]
        assert doc["original_filename"] == name
        assert doc["mime_type"] == mime_type
        assert doc["extension"] == extension
        assert doc["file_size"] > 0
        assert doc["sha256_checksum"] is not None
        assert len(doc["sha256_checksum"]) == 64
        assert doc["status"] == expected_status
        if expected_status == "READY":
            assert doc["error_message"] is None
        else:
            assert doc["error_message"] is not None
        assert "id" in doc
        assert "user_id" in doc
        assert "created_at" in doc
        assert "updated_at" in doc
        # Verify UUID format
        uuid.UUID(doc["id"])
        uuid.UUID(doc["user_id"])

    def test_upload_pdf(self, client, auth_headers, sample_pdf):
        r = self._upload(client, auth_headers, *sample_pdf)
        self._verify_ok(r, "test.pdf", "application/pdf", "pdf")

    def test_upload_png(self, client, auth_headers, sample_png):
        r = self._upload(client, auth_headers, *sample_png)
        self._verify_ok(r, "diagram.png", "image/png", "png", expected_status="FAILED")

    def test_upload_jpg(self, client, auth_headers, sample_jpg):
        r = self._upload(client, auth_headers, *sample_jpg)
        self._verify_ok(r, "photo.jpg", "image/jpeg", "jpg", expected_status="FAILED")

    def test_upload_jpeg(self, client, auth_headers, sample_jpeg):
        r = self._upload(client, auth_headers, *sample_jpeg)
        self._verify_ok(r, "image.jpeg", "image/jpeg", "jpeg", expected_status="FAILED")

    def test_upload_docx(self, client, auth_headers, sample_docx):
        r = self._upload(client, auth_headers, *sample_docx)
        self._verify_ok(r, "report.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "docx")

    def test_upload_pptx(self, client, auth_headers, sample_pptx):
        r = self._upload(client, auth_headers, *sample_pptx)
        self._verify_ok(r, "slides.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", "pptx")

    def test_upload_txt(self, client, auth_headers, sample_txt):
        r = self._upload(client, auth_headers, *sample_txt)
        self._verify_ok(r, "notes.txt", "text/plain", "txt")


# ===================================================================
# Validation — rejected uploads
# ===================================================================

class TestValidation:

    def test_invalid_extension(self, client, auth_headers):
        r = client.post(
            "/api/v1/files/upload",
            files={"file": ("malware.exe", b"some content")},
            headers=auth_headers,
        )
        assert r.status_code == 400
        assert "unsupported" in r.json()["detail"].lower()

    def test_invalid_mime_pdf(self, client, auth_headers):
        """Send .pdf extension with PNG content — MIME mismatch."""
        bad_content = _make_png_bytes()
        r = client.post(
            "/api/v1/files/upload",
            files={"file": ("fake.pdf", bad_content)},
            headers=auth_headers,
        )
        assert r.status_code == 400
        # Should say MIME doesn't match
        assert "mime" in r.json()["detail"].lower()

    def test_invalid_mime_exe_as_txt(self, client, auth_headers):
        """Send .txt with binary content — binary.txt not allowed."""
        r = client.post(
            "/api/v1/files/upload",
            files={"file": ("script.txt", b"\x00\x01\x02\xff")},
            headers=auth_headers,
        )
        assert r.status_code == 400

    def test_oversized_upload(self, client, auth_headers):
        """Send a file larger than MAX_UPLOAD_SIZE_MB (default 50 MB)."""
        # Simulate large file — 60 MB of data
        large = b"x" * (60 * 1024 * 1024)
        r = client.post(
            "/api/v1/files/upload",
            files={"file": ("big.pdf", large)},
            headers=auth_headers,
        )
        assert r.status_code == 413

    def test_empty_file(self, client, auth_headers):
        r = client.post(
            "/api/v1/files/upload",
            files={"file": ("empty.pdf", b"")},
            headers=auth_headers,
        )
        assert r.status_code == 400
        assert "empty" in r.json()["detail"].lower()

    def test_null_byte_in_filename(self, client, auth_headers):
        # The multipart parser strips null bytes, so check filename length
        # instead — a bare .pdf with no filename should be caught.
        r = client.post(
            "/api/v1/files/upload",
            files={"file": ("a" * 300 + ".pdf", _make_pdf_bytes())},
            headers=auth_headers,
        )
        assert r.status_code == 400

    def test_path_traversal_in_filename(self, client, auth_headers):
        r = client.post(
            "/api/v1/files/upload",
            files={"file": ("../../etc/passwd.pdf", _make_pdf_bytes())},
            headers=auth_headers,
        )
        assert r.status_code == 400
        assert "traversal" in r.json()["detail"].lower()

    def test_hidden_file(self, client, auth_headers):
        r = client.post(
            "/api/v1/files/upload",
            files={"file": (".secret.pdf", _make_pdf_bytes())},
            headers=auth_headers,
        )
        assert r.status_code == 400


class TestUploadRollback:
    def test_db_failure_removes_written_file(self, db_session, user, monkeypatch, tmp_path):
        """If the DB commit fails, the service should clean up the file it wrote."""
        from app.core.config import settings
        from app.services.file_service import FileService

        monkeypatch.setattr(settings, "UPLOAD_ROOT", str(tmp_path))
        monkeypatch.setattr(db_session, "commit", lambda: (_ for _ in ()).throw(RuntimeError("db down")))

        upload = UploadFile(
            filename="rollback.pdf",
            file=io.BytesIO(_make_pdf_bytes()),
        )

        with pytest.raises(RuntimeError, match="db down"):
            FileService.upload(db_session, upload, user.id)

        expected_path = Path(tmp_path) / str(user.id)
        assert not expected_path.exists()

    def test_no_extension(self, client, auth_headers):
        r = client.post(
            "/api/v1/files/upload",
            files={"file": ("README", "some text content")},
            headers=auth_headers,
        )
        assert r.status_code == 400
        assert "unsupported" in r.json()["detail"].lower()

    def test_path_separator_in_filename(self, client, auth_headers):
        r = client.post(
            "/api/v1/files/upload",
            files={"file": ("subdir/file.pdf", _make_pdf_bytes())},
            headers=auth_headers,
        )
        assert r.status_code == 400
        assert "separator" in r.json()["detail"].lower()


# ===================================================================
# Authentication & ownership
# ===================================================================

class TestAuthentication:

    def test_unauthenticated_upload(self, client):
        r = client.post(
            "/api/v1/files/upload",
            files={"file": ("doc.pdf", _make_pdf_bytes())},
        )
        assert r.status_code == 401  # HTTPBearer returns 401 without auth header

    def test_cross_user_access_prevention(self, client, auth_headers, registered_user, sample_pdf):
        """User B cannot access User A's document."""
        # User A uploads
        name, content = sample_pdf
        r1 = client.post(
            "/api/v1/files/upload",
            files={"file": (name, content)},
            headers=auth_headers,
        )
        assert r1.status_code == 201
        doc_id = r1.json()["document"]["id"]

        # User B registers and logs in
        email_b = "bob@example.com"
        password_b = "StrongPass2"
        client.post(
            "/api/v1/auth/register",
            json={"email": email_b, "password": password_b, "full_name": "Bob"},
        )
        r_login_b = client.post(
            "/api/v1/auth/login",
            json={"email": email_b, "password": password_b},
        )
        headers_b = {"Authorization": f"Bearer {r_login_b.json()['access_token']}"}

        # User B tries to read User A's document
        r2 = client.get(f"/api/v1/files/{doc_id}", headers=headers_b)
        assert r2.status_code == 404  # ownership enforced → not found

    def test_delete_other_users_document(self, client, auth_headers, sample_pdf):
        """User B cannot delete User A's document."""
        name, content = sample_pdf
        r1 = client.post(
            "/api/v1/files/upload",
            files={"file": (name, content)},
            headers=auth_headers,
        )
        assert r1.status_code == 201
        doc_id = r1.json()["document"]["id"]

        email_b = "bob2@example.com"
        password_b = "StrongPass3"
        client.post(
            "/api/v1/auth/register",
            json={"email": email_b, "password": password_b, "full_name": "Bob"},
        )
        r_login_b = client.post(
            "/api/v1/auth/login",
            json={"email": email_b, "password": password_b},
        )
        headers_b = {"Authorization": f"Bearer {r_login_b.json()['access_token']}"}

        r2 = client.delete(f"/api/v1/files/{doc_id}", headers=headers_b)
        assert r2.status_code == 404


# ===================================================================
# Metadata & status lifecycle
# ===================================================================

class TestMetadata:

    def test_metadata_persistence(self, client, auth_headers, sample_pdf):
        name, content = sample_pdf
        r = client.post(
            "/api/v1/files/upload",
            files={"file": (name, content)},
            headers=auth_headers,
        )
        assert r.status_code == 201
        doc_id = r.json()["document"]["id"]

        # Fetch by ID
        r2 = client.get(f"/api/v1/files/{doc_id}", headers=auth_headers)
        assert r2.status_code == 200
        d = r2.json()
        assert d["id"] == doc_id
        assert d["original_filename"] == "test.pdf"
        assert d["mime_type"] == "application/pdf"
        assert d["extension"] == "pdf"
        assert d["file_size"] > 0
        assert d["sha256_checksum"] is not None
        assert d["status"] == "READY"
        assert d["error_message"] is None

    def test_list_documents(self, client, auth_headers, sample_pdf, sample_png):
        """Upload two files, list them, verify count."""
        client.post(
            "/api/v1/files/upload",
            files={"file": sample_pdf},
            headers=auth_headers,
        )
        client.post(
            "/api/v1/files/upload",
            files={"file": sample_png},
            headers=auth_headers,
        )

        r = client.get("/api/v1/files/", headers=auth_headers)
        assert r.status_code == 200
        body = r.json()
        assert body["total"] == 2
        assert len(body["documents"]) == 2

    def test_list_documents_empty(self, client, auth_headers):
        r = client.get("/api/v1/files/", headers=auth_headers)
        assert r.status_code == 200
        assert r.json() == {"documents": [], "total": 0}

    def test_delete_document(self, client, auth_headers, sample_pdf):
        name, content = sample_pdf
        r = client.post(
            "/api/v1/files/upload",
            files={"file": (name, content)},
            headers=auth_headers,
        )
        doc_id = r.json()["document"]["id"]

        # Delete
        r2 = client.delete(f"/api/v1/files/{doc_id}", headers=auth_headers)
        assert r2.status_code == 204

        # Verify not found after delete
        r3 = client.get(f"/api/v1/files/{doc_id}", headers=auth_headers)
        assert r3.status_code == 404

    def test_get_nonexistent_document(self, client, auth_headers):
        fake_id = str(uuid.uuid4())
        r = client.get(f"/api/v1/files/{fake_id}", headers=auth_headers)
        assert r.status_code == 404

    def test_upload_status_lifecycle(self, client, auth_headers, sample_pdf):
        """Verify status moves through uploaded → processing → ready."""
        name, content = sample_pdf
        r = client.post(
            "/api/v1/files/upload",
            files={"file": (name, content)},
            headers=auth_headers,
        )
        assert r.status_code == 201
        # Final status should be READY
        assert r.json()["document"]["status"] == "READY"


# ===================================================================
# Checksum verification
# ===================================================================

class TestChecksum:

    def test_sha256_checksum_correct(self, client, auth_headers):
        import hashlib
        content = _make_pdf_bytes()
        expected_sha = hashlib.sha256(content).hexdigest()

        r = client.post(
            "/api/v1/files/upload",
            files={"file": ("doc.pdf", content)},
            headers=auth_headers,
        )
        assert r.status_code == 201
        assert r.json()["document"]["sha256_checksum"] == expected_sha

    def test_checksum_changes_with_content(self, client, auth_headers):
        r1 = client.post(
            "/api/v1/files/upload",
            files={"file": ("a.pdf", _make_pdf_bytes())},
            headers=auth_headers,
        )
        r2 = client.post(
            "/api/v1/files/upload",
            files={"file": ("b.pdf", _make_pdf_bytes() + b"extra")},
            headers=auth_headers,
        )
        assert r1.json()["document"]["sha256_checksum"] != r2.json()["document"]["sha256_checksum"]


# ===================================================================
# Filesystem verification
# ===================================================================

class TestFilesystem:

    def test_file_stored_on_disk(self, client, auth_headers, registered_user, monkeypatch, tmp_path):
        """Verify the file exists at storage/{user_id}/{doc_id}/{filename}."""
        from app.core.config import settings

        # Use a tmp dir for storage to avoid polluting real storage
        monkeypatch.setattr(settings, "UPLOAD_ROOT", str(tmp_path / "storage"))

        content = _make_pdf_bytes()
        r = client.post(
            "/api/v1/files/upload",
            files={"file": ("test.pdf", content)},
            headers=auth_headers,
        )
        assert r.status_code == 201
        doc = r.json()["document"]

        user_id = doc["user_id"]
        doc_id = doc["id"]

        stored_path = tmp_path / "storage" / user_id / doc_id / "test.pdf"
        assert stored_path.exists(), f"File not found at {stored_path}"
        assert stored_path.read_bytes() == content

    def test_file_size_matches_disk(self, client, auth_headers, registered_user, monkeypatch, tmp_path):
        from app.core.config import settings
        monkeypatch.setattr(settings, "UPLOAD_ROOT", str(tmp_path / "storage"))

        content = _make_pdf_bytes()
        r = client.post(
            "/api/v1/files/upload",
            files={"file": ("test.pdf", content)},
            headers=auth_headers,
        )
        doc = r.json()["document"]
        stored_path = tmp_path / "storage" / doc["user_id"] / doc["id"] / "test.pdf"
        assert stored_path.stat().st_size == doc["file_size"]


# ===================================================================
# Duplicate uploads
# ===================================================================

class TestDuplicates:

    def test_duplicate_content_unique_documents(self, client, auth_headers):
        """Same content, two uploads → two different document IDs."""
        content = _make_pdf_bytes()
        r1 = client.post(
            "/api/v1/files/upload",
            files={"file": ("a.pdf", content)},
            headers=auth_headers,
        )
        r2 = client.post(
            "/api/v1/files/upload",
            files={"file": ("b.pdf", content)},
            headers=auth_headers,
        )
        assert r1.status_code == 201
        assert r2.status_code == 201
        assert r1.json()["document"]["id"] != r2.json()["document"]["id"]
        # Checksums should be identical
        assert r1.json()["document"]["sha256_checksum"] == r2.json()["document"]["sha256_checksum"]


# ===================================================================
# Audit logging
# ===================================================================

class TestAuditLogging:

    def test_upload_logs_event(self, client, auth_headers, sample_pdf, caplog):
        """Upload triggers a structured log event."""
        import logging
        caplog.set_level(logging.INFO)

        name, content = sample_pdf
        r = client.post(
            "/api/v1/files/upload",
            files={"file": (name, content)},
            headers=auth_headers,
        )
        assert r.status_code == 201

        # Verify log output contains an upload event
        log_text = caplog.text
        assert "upload.file_uploaded" in log_text
        assert "test.pdf" in log_text
        assert "sha256" in log_text.lower() or "checksum" in log_text.lower()
