"""
Helpers for creating valid test documents with real extractable content.

Each function returns (filename, bytes) suitable for use with the
test upload endpoint. The documents contain real, extractable text
so the processing pipeline can process them.
"""

import io
from typing import Tuple


def make_pdf_bytes(text: str = "Hello world\nTest PDF content.\n") -> bytes:
    """Create a valid PDF with text content using fpdf2."""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    for line in text.split("\n"):
        if line.strip():
            pdf.set_font("Helvetica", size=12)
            # fpdf2 handles multi-byte text; for simple ASCII this works
            pdf.cell(0, 10, text=line.strip(), new_x="LMARGIN", new_y="NEXT")
    pdf_bytes = pdf.output()
    return pdf_bytes if isinstance(pdf_bytes, bytes) else bytes(pdf_bytes)


def make_docx_bytes(text: str = "Hello world\nTest DOCX content.\n") -> bytes:
    """Create a valid DOCX as bytes with text content."""
    from docx import Document

    doc = Document()
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line.strip())

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def make_pptx_bytes(text: str = "Hello world\nTest PPTX content.\n") -> bytes:
    """Create a valid PPTX as bytes with text content."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    added = False
    for line in text.split("\n"):
        if line.strip():
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.text = line.strip()
            added = True

    if not added:
        # Fallback: add at least one slide with placeholder text
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tf = txBox.text_frame
        tf.text = "Slide content"

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def make_txt_bytes(text: str = "Hello world\nTest TXT content.\n") -> bytes:
    """Create TXT bytes."""
    return text.encode("utf-8")
